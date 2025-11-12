import asyncio
import functools
import logging
import textwrap
import time
import rag_setup
from schemas import ChatRequest, DocumentRequest, TaskRequest
from typing import Optional, Tuple
from config import settings, detect_provider_from_key  # Fixed: removed 'app.' prefix
from fastapi import UploadFile, HTTPException
import json
import xml.etree.ElementTree as ET
from striprtf.striprtf import rtf_to_text
import markdown

try:
    import fitz  # PyMuPDF
except ImportError:
    logging.error("PyMuPDF is not installed. PDF processing will not work. Please run 'pip install pymupdf'")
    fitz = None

try:
    import docx  # python-docx for Word documents
except ImportError:
    logging.error(
        "python-docx is not installed. Word document processing will not work. Please run 'pip install python-docx'")
    docx = None

try:
    from pptx import Presentation  # python-pptx for PowerPoint
except ImportError:
    logging.error(
        "python-pptx is not installed. PowerPoint processing will not work. Please run 'pip install python-pptx'")
    Presentation = None

try:
    import pandas as pd  # For Excel and CSV files
except ImportError:
    logging.error(
        "pandas is not installed. Excel/CSV processing will not work. Please run 'pip install pandas openpyxl'")
    pd = None

try:
    from bs4 import BeautifulSoup  # For HTML parsing
except ImportError:
    logging.error(
        "BeautifulSoup is not installed. HTML processing will not work. Please run 'pip install beautifulsoup4'")
    BeautifulSoup = None

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%H:%M:%S'
)
logger = logging.getLogger("rag-service")

# A simple cache to store recent responses to avoid redundant API calls for the same query.
# The cache stores a tuple of (timestamp, response).
_response_cache = {}
CACHE_EXPIRATION_SECONDS = 600  # 10 minutes


def create_llm_instance(api_key: str, provider: Optional[str] = None):
    """Create a new LLM instance with the provided API key (OpenRouter or OpenAI)."""
    return rag_setup.create_llm(api_key=api_key, provider=provider)


async def test_api_key(api_key: str, provider: Optional[str] = None) -> dict:
    """Test if the provided API key is valid (OpenRouter or OpenAI)."""
    logger.info(f"🔍 Testing API key: {api_key[:10]}...")

    try:
        # Validate API key format first
        if not api_key or not api_key.strip():
            logger.error("❌ API key is empty")
            return {
                "valid": False,
                "message": "API key cannot be empty",
                "model_info": None
            }

        # Auto-detect provider if not specified
        if provider is None:
            provider = detect_provider_from_key(api_key)
            logger.info(f"🔍 Auto-detected provider: {provider}")

        # Validate based on provider
        if provider == "openrouter":
            if not api_key.startswith('sk-or-'):
                logger.error("❌ API key has incorrect format for OpenRouter")
                return {
                    "valid": False,
                    "message": "OpenRouter API keys should start with 'sk-or-'",
                    "model_info": None
                }
        elif provider == "openai":
            if not api_key.startswith('sk-'):
                logger.error("❌ API key has incorrect format for OpenAI")
                return {
                    "valid": False,
                    "message": "OpenAI API keys should start with 'sk-'",
                    "model_info": None
                }
        else:
            logger.error("❌ Unknown provider")
            return {
                "valid": False,
                "message": f"Unknown provider: {provider}. Please use OpenRouter or OpenAI.",
                "model_info": None
            }

        if len(api_key) < 40:
            logger.error("❌ API key is too short")
            return {
                "valid": False,
                "message": "API key appears to be too short",
                "model_info": None
            }

        # Create a temporary LLM instance
        test_llm = create_llm_instance(api_key, provider)

        # Test with a minimal prompt to avoid quota usage
        if provider == "openai":
            # OpenAI uses the SDK, so we test differently
            try:
                test_content = test_llm.generate_content("Hi", max_tokens=5)
                if test_content and not test_content.startswith("❌"):
                    logger.info("✅ OpenAI API key test successful")
                    return {
                        "valid": True,
                        "message": "OpenAI API key is valid and working!",
                        "model_info": {"model": settings.OPENAI_MODEL, "provider": "openai"}
                    }
                else:
                    return {
                        "valid": False,
                        "message": test_content or "API key test failed",
                        "model_info": None
                    }
            except Exception as e:
                error_msg = str(e)
                if "401" in error_msg or "Incorrect API key" in error_msg:
                    return {
                        "valid": False,
                        "message": "Invalid OpenAI API key: Authentication failed",
                        "model_info": None
                    }
                else:
                    return {
                        "valid": False,
                        "message": f"OpenAI API key test failed: {error_msg}",
                        "model_info": None
                    }

        # For OpenRouter, use the existing logic
        test_response = test_llm._make_api_request("Hi", max_tokens=1)

        # Check for explicit errors first
        if "error" in test_response:
            error_msg = test_response["error"]
            logger.error(f"❌ API key test failed: {error_msg}")

            # Parse specific error types
            if "401" in str(error_msg) or "403" in str(error_msg) or "Unauthorized" in str(error_msg):
                return {
                    "valid": False,
                    "message": "Invalid API key: Authentication failed",
                    "model_info": None
                }
            elif "429" in str(error_msg):
                return {
                    "valid": False,
                    "message": "API key valid but rate limited. Try again in a moment.",
                    "model_info": None
                }
            elif "insufficient" in str(error_msg).lower() or "quota" in str(error_msg).lower():
                # API key is valid but no credits - this is still a valid key
                logger.info("✅ API key is valid but has no credits")
                return {
                    "valid": True,
                    "message": "API key is valid but has insufficient credits",
                    "model_info": {"model": "unknown", "usage": {}, "credits": "insufficient"}
                }
            else:
                return {
                    "valid": False,
                    "message": f"API key test failed: {error_msg}",
                    "model_info": None
                }

        # Check for successful response with proper structure
        if "choices" in test_response and test_response["choices"]:
            choice = test_response["choices"][0]
            if "message" in choice and "content" in choice["message"]:
                logger.info("✅ API key test successful")
                model_info = {
                    "model": test_response.get("model", "unknown"),
                    "usage": test_response.get("usage", {}),
                    "credits": "available"
                }
                return {
                    "valid": True,
                    "message": "API key is valid and working!",
                    "model_info": model_info
                }

        # If we get here, the response format is unexpected
        logger.error(f"❌ API key test failed: Unexpected response format - {test_response}")
        return {
            "valid": False,
            "message": "API key test failed: Unexpected response format from OpenRouter",
            "model_info": None
        }

    except Exception as e:
        logger.error(f"❌ API key test failed with exception: {str(e)}")
        error_msg = str(e)

        # Parse common error patterns
        if "401" in error_msg or "403" in error_msg or "Unauthorized" in error_msg:
            return {
                "valid": False,
                "message": "Invalid API key: Authentication failed",
                "model_info": None
            }
        elif "timeout" in error_msg.lower():
            return {
                "valid": False,
                "message": "API key test timed out. Please try again.",
                "model_info": None
            }
        elif "connection" in error_msg.lower():
            return {
                "valid": False,
                "message": "Network connection error. Please check your internet connection.",
                "model_info": None
            }
        else:
            return {
                "valid": False,
                "message": f"API key test failed: {error_msg}",
                "model_info": None
            }


async def process_and_index_file(file: UploadFile) -> Tuple[int, str]:
    """
    Processes an uploaded file, extracts text, calls the indexing function,
    and returns the number of documents and the extracted text.
    Supports: .txt, .pdf, .docx, .pptx, .xlsx, .csv, .json, .xml, .html, .md, .rtf
    """
    logger.info(f"📄 Processing file '{file.filename}' with content type '{file.content_type}'")

    # Read file content
    file_content = await file.read()
    text = ""
    file_extension = file.filename.lower().split('.')[-1]

    try:
        if file_extension == "txt":
            text = await _process_txt_file(file_content)

        elif file_extension == "pdf":
            text = await _process_pdf_file(file_content)

        elif file_extension == "docx":
            text = await _process_docx_file(file_content)

        elif file_extension in ["ppt", "pptx"]:
            text = await _process_pptx_file(file_content)

        elif file_extension in ["xls", "xlsx"]:
            text = await _process_excel_file(file_content, file.filename)

        elif file_extension == "csv":
            text = await _process_csv_file(file_content)

        elif file_extension == "json":
            text = await _process_json_file(file_content)

        elif file_extension == "xml":
            text = await _process_xml_file(file_content)

        elif file_extension in ["html", "htm"]:
            text = await _process_html_file(file_content)

        elif file_extension in ["md", "markdown"]:
            text = await _process_markdown_file(file_content)

        elif file_extension == "rtf":
            text = await _process_rtf_file(file_content)

        else:
            supported_extensions = ['.txt', '.pdf', '.docx', '.pptx', '.xlsx', '.csv', '.json', '.xml', '.html', '.md',
                                    '.rtf']
            logger.error(f"❌ Unsupported file type: {file.filename}")
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Please upload one of: {', '.join(supported_extensions)}"
            )

    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        logger.error(f"❌ Failed to process {file_extension.upper()} file: {e}")
        raise HTTPException(status_code=400, detail=f"Could not process {file_extension.upper()} file: {str(e)}")

    # Validate extracted text
    if not text or not text.strip():
        logger.error("❌ Extracted text is empty or whitespace only")
        raise HTTPException(status_code=400,
                            detail="Extracted text is empty. The file might be empty, corrupted, or unreadable.")

    # Clean up the text
    text = text.strip()

    # Log processing stats
    word_count = len(text.split())
    logger.info(f"📊 Text processing complete: {len(text)} characters, {word_count} words")

    # Index the extracted text using existing logic
    try:
        doc_request = DocumentRequest(context=text)
        docs_added = index_document(doc_request)
        logger.info(f"✅ Successfully indexed {docs_added} document chunks from file")

        return docs_added, text

    except Exception as e:
        logger.error(f"❌ Failed to index extracted text: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to index extracted text: {str(e)}")


# Individual file processing functions
async def _process_txt_file(file_content: bytes) -> str:
    """Process .txt files with multiple encoding support."""
    try:
        # Try UTF-8 first, then fallback to other encodings
        try:
            text = file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Try other common encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    text = file_content.decode(encoding)
                    logger.info(f"✅ Decoded text file using {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise UnicodeDecodeError("Unable to decode file with any common encoding")

        logger.info(f"✅ Extracted {len(text)} characters from .txt file")
        return text

    except UnicodeDecodeError as e:
        logger.error(f"❌ Could not decode .txt file: {e}")
        raise HTTPException(status_code=400,
                            detail="Could not decode .txt file. Please ensure it uses UTF-8, Latin-1, or CP1252 encoding.")


async def _process_pdf_file(file_content: bytes) -> str:
    """Process .pdf files using PyMuPDF."""
    if fitz is None:
        logger.error("❌ PyMuPDF not available for PDF processing")
        raise HTTPException(status_code=501, detail="PDF processing is not available. PyMuPDF is not installed.")

    logger.info("📖 Opening PDF document...")
    doc = fitz.open(stream=file_content, filetype="pdf")

    try:
        text_parts = []
        page_count = len(doc)
        logger.info(f"📑 PDF has {page_count} pages")

        for page_num in range(page_count):
            try:
                page = doc[page_num]
                page_text = page.get_text()

                if page_text and page_text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}")
                    logger.info(f"📄 Extracted text from page {page_num + 1}: {len(page_text)} characters")
                else:
                    logger.info(f"📄 Page {page_num + 1} is empty or contains no extractable text")

            except Exception as page_error:
                logger.warning(f"⚠️  Could not extract text from page {page_num + 1}: {page_error}")
                continue

        text = "\n\n".join(text_parts)
        logger.info(f"✅ Extracted text from {len(text_parts)} pages of the PDF file ({len(text)} characters)")
        return text

    finally:
        doc.close()
        logger.info("📕 PDF document closed successfully")


async def _process_docx_file(file_content: bytes) -> str:
    """Process .docx files using python-docx."""
    if docx is None:
        raise HTTPException(status_code=501,
                            detail="Word document processing is not available. python-docx is not installed.")

    from io import BytesIO
    doc = docx.Document(BytesIO(file_content))

    text_parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text.strip())

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                if cell.text.strip():
                    row_text.append(cell.text.strip())
            if row_text:
                text_parts.append(" | ".join(row_text))

    text = "\n\n".join(text_parts)
    logger.info(f"✅ Extracted {len(text)} characters from Word document")
    return text


async def _process_pptx_file(file_content: bytes) -> str:
    """Process .pptx files using python-pptx."""
    if Presentation is None:
        raise HTTPException(status_code=501,
                            detail="PowerPoint processing is not available. python-pptx is not installed.")

    from io import BytesIO
    prs = Presentation(BytesIO(file_content))

    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if len(slide_text) > 1:  # More than just the slide header
            text_parts.append("\n".join(slide_text))

    text = "\n\n".join(text_parts)
    logger.info(f"✅ Extracted text from {len(prs.slides)} PowerPoint slides ({len(text)} characters)")
    return text


async def _process_excel_file(file_content: bytes, filename: str) -> str:
    """Process .xlsx/.xls files using pandas."""
    if pd is None:
        raise HTTPException(status_code=501, detail="Excel processing is not available. pandas is not installed.")

    from io import BytesIO

    try:
        # Read all sheets
        excel_file = pd.ExcelFile(BytesIO(file_content))
        text_parts = [f"Excel File: {filename}"]

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)

            if not df.empty:
                text_parts.append(f"\n--- Sheet: {sheet_name} ---")

                # Convert DataFrame to readable text
                # Include column headers
                text_parts.append("Columns: " + " | ".join(str(col) for col in df.columns))

                # Add row data (limit to first 100 rows to avoid huge files)
                for idx, row in df.head(100).iterrows():
                    row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
                    if row_text.strip():
                        text_parts.append(row_text)

                if len(df) > 100:
                    text_parts.append(f"... and {len(df) - 100} more rows")

        text = "\n".join(text_parts)
        logger.info(
            f"✅ Extracted data from Excel file with {len(excel_file.sheet_names)} sheets ({len(text)} characters)")
        return text

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process Excel file: {str(e)}")


async def _process_csv_file(file_content: bytes) -> str:
    """Process .csv files using pandas."""
    if pd is None:
        raise HTTPException(status_code=501, detail="CSV processing is not available. pandas is not installed.")

    from io import StringIO

    try:
        # Try different encodings for CSV
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                csv_text = file_content.decode(encoding)
                df = pd.read_csv(StringIO(csv_text))
                break
            except (UnicodeDecodeError, pd.errors.EmptyDataError):
                continue
        else:
            raise ValueError("Could not decode CSV file with any common encoding")

        if df.empty:
            raise ValueError("CSV file is empty")

        text_parts = ["CSV Data:"]
        text_parts.append("Columns: " + " | ".join(str(col) for col in df.columns))

        # Add row data (limit to first 200 rows)
        for idx, row in df.head(200).iterrows():
            row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
            if row_text.strip():
                text_parts.append(row_text)

        if len(df) > 200:
            text_parts.append(f"... and {len(df) - 200} more rows")

        text = "\n".join(text_parts)
        logger.info(f"✅ Extracted data from CSV file with {len(df)} rows ({len(text)} characters)")
        return text

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process CSV file: {str(e)}")


async def _process_json_file(file_content: bytes) -> str:
    """Process .json files."""
    try:
        json_text = file_content.decode('utf-8')
        data = json.loads(json_text)

        # Convert JSON to readable text format
        def json_to_text(obj, indent=0):
            lines = []
            prefix = "  " * indent

            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, (dict, list)):
                        lines.append(f"{prefix}{key}:")
                        lines.extend(json_to_text(value, indent + 1))
                    else:
                        lines.append(f"{prefix}{key}: {value}")
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    if isinstance(item, (dict, list)):
                        lines.append(f"{prefix}[{i}]:")
                        lines.extend(json_to_text(item, indent + 1))
                    else:
                        lines.append(f"{prefix}[{i}]: {item}")
            else:
                lines.append(f"{prefix}{obj}")

            return lines

        text_lines = ["JSON Data:"] + json_to_text(data)
        text = "\n".join(text_lines)

        logger.info(f"✅ Extracted data from JSON file ({len(text)} characters)")
        return text

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"Invalid JSON file: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process JSON file: {str(e)}")


async def _process_xml_file(file_content: bytes) -> str:
    """Process .xml files."""
    try:
        xml_text = file_content.decode('utf-8')
        root = ET.fromstring(xml_text)

        def xml_to_text(element, indent=0):
            lines = []
            prefix = "  " * indent

            # Add element name and attributes
            if element.attrib:
                attrs = " ".join(f'{k}="{v}"' for k, v in element.attrib.items())
                lines.append(f"{prefix}{element.tag} ({attrs}):")
            else:
                lines.append(f"{prefix}{element.tag}:")

            # Add text content
            if element.text and element.text.strip():
                lines.append(f"{prefix}  {element.text.strip()}")

            # Add child elements
            for child in element:
                lines.extend(xml_to_text(child, indent + 1))

            return lines

        text_lines = ["XML Data:"] + xml_to_text(root)
        text = "\n".join(text_lines)

        logger.info(f"✅ Extracted data from XML file ({len(text)} characters)")
        return text

    except ET.ParseError as e:
        raise HTTPException(status_code=400, detail=f"Invalid XML file: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process XML file: {str(e)}")


async def _process_html_file(file_content: bytes) -> str:
    """Process .html files using BeautifulSoup."""
    if BeautifulSoup is None:
        raise HTTPException(status_code=501, detail="HTML processing is not available. BeautifulSoup is not installed.")

    try:
        html_text = file_content.decode('utf-8')
        soup = BeautifulSoup(html_text, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text content
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        logger.info(f"✅ Extracted text from HTML file ({len(text)} characters)")
        return text

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process HTML file: {str(e)}")


async def _process_markdown_file(file_content: bytes) -> str:
    """Process .md files."""
    try:
        md_text = file_content.decode('utf-8')

        # Convert markdown to HTML then to plain text for better readability
        html = markdown.markdown(md_text)
        if BeautifulSoup:
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
        else:
            # Fallback: use raw markdown
            text = md_text

        logger.info(f"✅ Extracted text from Markdown file ({len(text)} characters)")
        return text

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process Markdown file: {str(e)}")


async def _process_rtf_file(file_content: bytes) -> str:
    """Process .rtf files."""
    try:
        rtf_text = file_content.decode('utf-8')
        text = rtf_to_text(rtf_text)

        logger.info(f"✅ Extracted text from RTF file ({len(text)} characters)")
        return text

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not process RTF file: {str(e)}")


def _create_overlapping_chunks(text: str, chunk_size: int = 500, overlap: int = 100) -> list:
    """
    Create overlapping text chunks for better context preservation.
    """
    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = start + chunk_size

        if end < text_length:
            # Look for sentence boundaries
            sentence_end = max(
                text.rfind('.', start, end),
                text.rfind('!', start, end),
                text.rfind('?', start, end)
            )

            if sentence_end > start + chunk_size // 2:
                end = sentence_end + 1
            else:
                space_pos = text.rfind(' ', start, end)
                if space_pos > start + chunk_size // 2:
                    end = space_pos

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        # Ensure we always move forward to avoid infinite loops
        next_start = end - overlap if end < text_length else text_length
        start = max(next_start, start + 1)

    return chunks


def index_document(request_data: DocumentRequest) -> int:
    logger.info("=" * 80)
    logger.info("📚 STARTING ENHANCED DOCUMENT INDEXING PROCESS")
    logger.info("=" * 80)

    # Log the incoming context
    context_preview = request_data.context[:200] + "..." if len(request_data.context) > 200 else request_data.context
    logger.info(f"📝 CONTEXT TO INDEX (length: {len(request_data.context)} chars):")
    logger.info(f"   Preview: {context_preview}")
    logger.info("-" * 60)

    try:
        # Step 1: Clear any existing documents properly
        existing_data = rag_setup.collection.get()
        existing_ids = existing_data.get("ids", [])
        if existing_ids:
            rag_setup.collection.delete(ids=existing_ids)
            logger.info(f"🗑️  Cleared {len(existing_ids)} existing documents from vector collection.")
        else:
            logger.info("📂 No existing documents to clear.")

        # Step 2: Enhanced chunking with overlap for better context preservation
        logger.info("✂️  Creating overlapping chunks for better context continuity...")
        text_chunks = _create_overlapping_chunks(
            request_data.context,
            chunk_size=500,  # Optimal size for retrieval
            overlap=100  # 20% overlap for context preservation
        )

        if not text_chunks:
            logger.warning("⚠️  No text chunks were generated.")
            return 0

        logger.info(f"✅ Document split into {len(text_chunks)} overlapping chunks")

        # Log chunk statistics
        avg_chunk_size = sum(len(chunk) for chunk in text_chunks) / len(text_chunks)
        logger.info(f"📊 Average chunk size: {avg_chunk_size:.0f} characters")

        # Step 3: Add chunks to ChromaDB with enhanced metadata
        timestamp = int(time.time())
        chunk_ids = [f"doc_chunk_{i}_{timestamp}" for i in range(len(text_chunks))]
        logger.info(f"💾 Adding {len(chunk_ids)} chunks to ChromaDB with enhanced metadata...")

        # Enhanced metadata for better retrieval
        metadatas = [
            {
                "chunk_index": i,
                "timestamp": timestamp,
                "chunk_length": len(chunk),
                "position": "start" if i == 0 else "end" if i == len(text_chunks) - 1 else "middle",
                "total_chunks": len(text_chunks)
            }
            for i, chunk in enumerate(text_chunks)
        ]

        rag_setup.collection.add(
            documents=text_chunks,
            ids=chunk_ids,
            metadatas=metadatas
        )

        logger.info("✅ ENHANCED DOCUMENT INDEXING COMPLETED SUCCESSFULLY")
        logger.info(f"📊 Total chunks indexed: {len(text_chunks)}")
        logger.info(f"🔗 Chunks have 100-char overlap for better context continuity")
        logger.info("=" * 80)

        return len(text_chunks)

    except Exception as e:
        logger.error(f"❌ Error during indexing: {str(e)}", exc_info=True)
        raise


def clear_index():
    """Clears all documents from the vector database."""
    logger.info("🗑️  Clearing vector index...")
    try:
        # Get all existing documents and delete them
        existing_data = rag_setup.collection.get()
        existing_ids = existing_data.get("ids", [])
        if existing_ids:
            rag_setup.collection.delete(ids=existing_ids)
            logger.info(f"✅ Successfully cleared {len(existing_ids)} documents from the vector index.")
        else:
            logger.info("✅ Vector index was already empty.")
    except Exception as e:
        logger.error(f"❌ Error clearing vector index: {e}")
        raise


def _expand_query(query: str) -> list:
    """
    Expand the query with synonyms and related terms for better retrieval.
    Returns list of query variations.
    """
    queries = [query]

    # Add query without question words for better matching
    question_words = ['what', 'where', 'when', 'who', 'why', 'how', 'is', 'are', 'can', 'do', 'does']
    words = query.lower().split()
    filtered_words = [w for w in words if w not in question_words and len(w) > 2]

    if len(filtered_words) >= 2:
        # Create a version with just the key terms
        key_terms_query = ' '.join(filtered_words)
        if key_terms_query != query.lower():
            queries.append(key_terms_query)

    return queries[:2]  # Limit to 2 variations to avoid too many retrievals


def _deduplicate_and_rank_chunks(chunks_list: list, metadatas_list: list) -> tuple:
    """
    Deduplicate chunks and rank them by relevance.
    Returns (unique_chunks, unique_metadatas)
    """
    seen = set()
    unique_chunks = []
    unique_metadatas = []

    for chunks, metadatas in zip(chunks_list, metadatas_list):
        for chunk, metadata in zip(chunks, metadatas):
            # Use first 100 chars as fingerprint
            fingerprint = chunk[:100]
            if fingerprint not in seen:
                seen.add(fingerprint)
                unique_chunks.append(chunk)
                unique_metadatas.append(metadata)

    return unique_chunks, unique_metadatas


async def get_rag_response(request_data: ChatRequest, api_key: Optional[str] = None) -> str:
    """
    Enhanced RAG pipeline with conversation history, query expansion, and better retrieval.
    """
    start_total = time.time()

    logger.info("=" * 80)
    logger.info("🤖 STARTING ENHANCED RAG PIPELINE")
    logger.info("=" * 80)
    logger.info(f"❓ USER PROMPT: '{request_data.prompt}'")
    logger.info(f"📏 Prompt length: {len(request_data.prompt)} characters")
    logger.info(f"💬 Conversation history: {len(request_data.conversation_history or [])} messages")
    logger.info(f"🔑 Using custom API key: {'Yes' if api_key else 'No'}")
    logger.info("-" * 60)

    try:
        # Step 1: Build cache key including conversation context
        history_hash = ""
        if request_data.conversation_history:
            # Create a hash of recent conversation for cache key
            recent_msgs = request_data.conversation_history[-3:]  # Last 3 messages
            history_hash = str(hash("".join([m.content[:50] for m in recent_msgs])))

        cache_key = f"{api_key or 'default'}:{history_hash}:{request_data.prompt}"
        cached_response = _get_cached_response(cache_key)
        if cached_response:
            logger.info("💾 CACHE HIT! Returning cached response.")
            return cached_response

        logger.info("🔍 Cache miss. Proceeding with enhanced RAG pipeline.")

        # Step 2: Check if the vector database has any content
        doc_count = rag_setup.collection.count()
        logger.info(f"📚 Vector DB contains {doc_count} documents")

        if doc_count == 0:
            logger.warning("⚠️  Vector DB is empty. Cannot answer query.")
            return "I don't have any specific context loaded right now. Please provide some context in the Knowledge Base and click 'Index Context' before asking questions. However, I'd be happy to help with general questions using my built-in knowledge!"

        # Step 3: Query expansion for better retrieval
        logger.info("🔍 Expanding query for better retrieval...")
        query_variations = _expand_query(request_data.prompt)
        logger.info(f"📝 Generated {len(query_variations)} query variations")

        # Step 4: Retrieve chunks for each query variation
        all_chunks = []
        all_metadatas = []

        for query_var in query_variations:
            logger.info(f"🔎 Retrieving chunks for: '{query_var[:50]}...'")
            retrieved = await _retrieve_chunks_async(
                query_var,
                n_results=settings.MAX_CHUNKS_RETRIEVE + 2  # Retrieve more for better ranking
            )

            if retrieved and retrieved.get('documents') and retrieved['documents'][0]:
                all_chunks.append(retrieved['documents'][0])
                all_metadatas.append(retrieved.get('metadatas', [[]])[0])

        if not all_chunks:
            logger.warning("❌ No relevant chunks found in the vector DB for this query.")
            return "I couldn't find specific information about that in the provided context. Let me help you with what I know from my general knowledge:\n\n" + await _generate_fallback_response(
                request_data.prompt, api_key)

        # Step 5: Deduplicate and rank chunks
        logger.info("🎯 Deduplicating and ranking retrieved chunks...")
        unique_chunks, unique_metadatas = _deduplicate_and_rank_chunks(all_chunks, all_metadatas)

        # Limit to best chunks
        final_chunks = unique_chunks[:settings.MAX_CHUNKS_RETRIEVE + 1]
        logger.info(f"📋 Using {len(final_chunks)} unique, ranked chunks for context")

        # Log chunk details
        for i, (chunk, meta) in enumerate(zip(final_chunks, unique_metadatas[:len(final_chunks)])):
            logger.info(f"   Chunk {i + 1}: {len(chunk)} chars, position: {meta.get('position', 'unknown')}")

        context_for_prompt = "\n\n---\n\n".join(final_chunks)

        # Limit context length to prevent timeouts
        max_context_length = settings.MAX_CONTEXT_LENGTH_CHAT
        if len(context_for_prompt) > max_context_length:
            logger.warning(f"⚠️  Context too long, truncating to {max_context_length}")
            context_for_prompt = context_for_prompt[:max_context_length] + "\n\n[... content truncated ...]"

        # Step 6: Build conversation history context
        history_context = ""
        if request_data.conversation_history and len(request_data.conversation_history) > 0:
            logger.info(f"💬 Including {len(request_data.conversation_history)} previous messages for context")
            history_messages = request_data.conversation_history[-6:]  # Last 6 messages (3 exchanges)
            history_parts = []
            for msg in history_messages:
                role_label = "User" if msg.role == "user" else "Assistant"
                # Truncate very long messages
                content = msg.content[:300] + "..." if len(msg.content) > 300 else msg.content
                history_parts.append(f"{role_label}: {content}")

            history_context = (
                "\n\nPREVIOUS CONVERSATION:\n"
                + "\n".join(history_parts)
                + "\n"
            )

        # Step 7: Construct enhanced prompt with conversation history
        full_prompt = (
            "You are an intelligent assistant with access to specific context information and conversation history. "
            "Your goal is to provide comprehensive, helpful answers that:\n"
            "• Take into account the previous conversation flow\n"
            "• Use the provided context as your PRIMARY source when relevant\n"
            "• Build naturally on previous exchanges\n"
            "• Provide coherent, contextually appropriate responses\n\n"

            "INSTRUCTIONS:\n"
            "• Reference previous conversation when relevant for continuity\n"
            "• Use the document context as your primary source for factual information\n"
            "• If the user asks follow-up questions, refer back to previous answers\n"
            "• Be natural and conversational - maintain the conversation flow\n"
            "• Provide detailed, well-structured responses\n"
            "• If information isn't in the context, acknowledge it and provide general knowledge help\n"
        )

        if history_context:
            full_prompt += history_context + "\n"

        full_prompt += (
            "DOCUMENT CONTEXT:\n"
            f"{context_for_prompt}\n\n"

            f"CURRENT USER QUESTION: {request_data.prompt}\n\n"

            "Please provide a comprehensive, contextually appropriate response:"
        )

        # Step 8: Generate the response using the LLM
        logger.info("🧠 Generating response with conversation context...")
        response_text = await _generate_response_async(full_prompt, api_key)

        # Step 9: Cache the newly generated response
        _cache_response(cache_key, response_text)
        logger.info("💾 Response cached for future use")

        total_time = time.time() - start_total
        logger.info(f"⏱️  Total processing time: {total_time:.2f}s")
        logger.info("✅ ENHANCED RAG PIPELINE COMPLETED SUCCESSFULLY")
        logger.info("=" * 80)

        return response_text

    except asyncio.TimeoutError:
        logger.error("⏱️  Request timed out during retrieval or generation.")
        return "The request timed out. Please try again or simplify your question."
    except Exception as e:
        logger.error(f"❌ An unexpected error occurred: {e}", exc_info=True)
        return f"An unexpected error occurred: {e}"


async def _generate_fallback_response(prompt: str, api_key: Optional[str] = None) -> str:
    """Generate a response using only the model's knowledge when no context is available."""
    fallback_prompt = (
        f"Please provide a helpful, comprehensive answer to this question using your knowledge:\n\n"
        f"Question: {prompt}\n\n"
        f"Answer:"
    )

    try:
        return await _generate_response_async(fallback_prompt, api_key)
    except Exception as e:
        logger.error(f"❌ Fallback response generation failed: {e}")
        return "I'm having trouble generating a response right now. Please try again or rephrase your question."


async def execute_task(request_data: TaskRequest, api_key: Optional[str] = None) -> str:
    """
    Executes a specific task on the given context.
    Uses provided API key or falls back to default.
    """
    start_total = time.time()

    logger.info("=" * 80)
    logger.info("🎯 STARTING TASK EXECUTION")
    logger.info("=" * 80)
    logger.info(f"📋 TASK TYPE: {request_data.task_type}")
    logger.info(f"❓ TASK PROMPT: '{request_data.prompt}'")
    logger.info(f"📏 Context length: {len(request_data.context)} characters")
    logger.info(f"🔑 Using custom API key: {'Yes' if api_key else 'No'}")
    logger.info("-" * 60)

    try:
        context = request_data.context
        if not context:
            logger.warning("⚠️  Context is empty for task execution")
            return "Context is empty. Please provide some text in the 'Knowledge Base' to perform a task."

        # Limit context length to prevent timeouts
        max_context_length = settings.MAX_CONTEXT_LENGTH_TASK
        if len(context) > max_context_length:
            logger.warning(f"⚠️  Context too long, truncating to {max_context_length}")
            context = context[:max_context_length] + "\n\n[... content truncated for performance ...]"

        # Construct optimized prompts based on the task type
        if request_data.task_type == "summarize":
            full_prompt = f"Summarize this text concisely:\n\n{context}"
        elif request_data.task_type == "plan":
            goal = request_data.prompt or 'Create a general action plan'
            full_prompt = f"Create an action plan based on this context.\n\nGoal: {goal}\n\nContext:\n{context}\n\nPlan:"
        elif request_data.task_type == "creative":
            style = request_data.prompt or 'Write a short creative piece'
            full_prompt = f"Write a creative piece inspired by this text.\n\nStyle/Request: {style}\n\nInspiration:\n{context}\n\nCreative piece:"
        else:
            logger.error(f"❌ Invalid task type: {request_data.task_type}")
            return "Invalid task type specified."

        # Generate the response
        logger.info("🧠 Generating task-based response from OpenRouter...")
        response_text = await _generate_response_async(full_prompt, api_key)

        total_time = time.time() - start_total
        logger.info(f"⏱️  Task execution time: {total_time:.2f}s")
        logger.info("✅ TASK EXECUTION COMPLETED SUCCESSFULLY")
        logger.info("=" * 80)

        return response_text

    except asyncio.TimeoutError:
        logger.error("⏱️  Request timed out during task execution.")
        return "The request timed out. Please try again."
    except Exception as e:
        logger.error(f"❌ An unexpected error occurred during task execution: {e}", exc_info=True)
        return f"An unexpected error occurred: {e}"


# --- ASYNC WRAPPERS & CACHE HELPERS ---

async def _retrieve_chunks_async(prompt: str, n_results: int = 2):
    """Asynchronously queries the ChromaDB collection."""
    logger.info(f"🔍 Querying ChromaDB for prompt: '{prompt}' (requesting {n_results} chunks)")
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(
        None,
        functools.partial(rag_setup.collection.query, query_texts=[prompt], n_results=n_results)
    )
    logger.info(f"📊 ChromaDB query returned {len(result.get('documents', [[]])[0])} chunks")
    return result


async def _generate_response_async(full_prompt: str, api_key: Optional[str] = None):
    """Asynchronously calls the LLM to generate content."""
    logger.info("🤖 Calling LLM for content generation...")
    logger.info(f"📏 Prompt length sent to LLM: {len(full_prompt)} characters")

    # Use custom API key if provided, otherwise use default
    if api_key:
        llm_instance = create_llm_instance(api_key)
        logger.info("🔑 Using user-provided API key")
    else:
        llm_instance = rag_setup.generation_model
        logger.info("🔧 Using default API key")

    loop = asyncio.get_event_loop()
    response = await loop.run_in_executor(
        None,
        llm_instance.generate_content,
        full_prompt
    )

    logger.info(f"✅ LLM response received (length: {len(response)} chars)")
    return response


def _get_cached_response(key: str):
    """Checks the cache for a valid (non-expired) entry."""
    if key in _response_cache:
        timestamp, response = _response_cache[key]
        if time.time() - timestamp < CACHE_EXPIRATION_SECONDS:
            logger.info(f"💾 Cache hit for key: '{key[:50]}...'")
            return response
        else:
            # Expired, remove from cache
            del _response_cache[key]
            logger.info(f"🗑️  Expired cache entry removed for key: '{key[:50]}...'")
    return None


def _cache_response(key: str, response: str):
    """Adds a response to the cache with the current timestamp."""
    _response_cache[key] = (time.time(), response)
    logger.info(f"💾 Response cached for key: '{key[:50]}...' (response length: {len(response)} chars)")