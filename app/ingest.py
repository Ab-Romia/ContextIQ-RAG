"""Document text extraction.

Turns an uploaded file into plain text plus a source label. Each format has a small
dedicated extractor; the dispatcher picks one by file extension. Extractors raise
`UnsupportedFile` or `ExtractionError` and stay free of any web-framework types, so
the same code runs from the API, the evaluation harness, or a unit test.
"""

from __future__ import annotations

import csv
import io
import json

import defusedxml.ElementTree as ET
from xml.etree.ElementTree import Element

TEXT_ENCODINGS = ("utf-8", "latin-1", "cp1252")


class UnsupportedFile(Exception):
    """Raised when a file extension has no extractor."""


class ExtractionError(Exception):
    """Raised when a supported file cannot be read."""


def _decode(content: bytes) -> str:
    for encoding in TEXT_ENCODINGS:
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ExtractionError("Could not decode file as UTF-8, Latin-1, or CP1252.")


def _extract_txt(content: bytes) -> str:
    return _decode(content)


def _extract_pdf(content: bytes) -> str:
    import fitz  # PyMuPDF

    parts: list[str] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            text = page.get_text().strip()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(content: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(content))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = [f"Slide {index}"]
        for shape in slide.shapes:
            if getattr(shape, "text", "").strip():
                lines.append(shape.text.strip())
        if len(lines) > 1:
            parts.append("\n".join(lines))
    return "\n\n".join(parts)


def _extract_excel(content: bytes) -> str:
    import pandas as pd

    parts: list[str] = []
    workbook = pd.ExcelFile(io.BytesIO(content))
    for sheet in workbook.sheet_names:
        frame = workbook.parse(sheet)
        if frame.empty:
            continue
        parts.append(f"Sheet: {sheet}")
        parts.append(frame.to_csv(index=False).strip())
    return "\n\n".join(parts)


def _extract_csv(content: bytes) -> str:
    text = _decode(content)
    rows = list(csv.reader(io.StringIO(text)))
    return "\n".join(" | ".join(cell for cell in row if cell) for row in rows if any(row))


def _extract_json(content: bytes) -> str:
    data = json.loads(_decode(content))
    return json.dumps(data, indent=2, ensure_ascii=False)


def _extract_xml(content: bytes) -> str:
    root = ET.fromstring(_decode(content))

    def walk(node: Element, depth: int = 0) -> list[str]:
        pad = "  " * depth
        lines = [f"{pad}{node.tag}: {node.text.strip()}" if node.text and node.text.strip() else f"{pad}{node.tag}"]
        for child in node:
            lines.extend(walk(child, depth + 1))
        return lines

    return "\n".join(walk(root))


def _extract_html(content: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(content), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    lines = [line.strip() for line in soup.get_text("\n").splitlines()]
    return "\n".join(line for line in lines if line)


def _extract_markdown(content: bytes) -> str:
    # Markdown is already readable text; keep the source markup so heading-aware
    # chunking can use the structure later.
    return _decode(content)


def _extract_rtf(content: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(_decode(content))


_EXTRACTORS = {
    "txt": _extract_txt,
    "text": _extract_txt,
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "ppt": _extract_pptx,
    "xlsx": _extract_excel,
    "xls": _extract_excel,
    "csv": _extract_csv,
    "json": _extract_json,
    "xml": _extract_xml,
    "html": _extract_html,
    "htm": _extract_html,
    "md": _extract_markdown,
    "markdown": _extract_markdown,
    "rtf": _extract_rtf,
}

SUPPORTED_EXTENSIONS = sorted(_EXTRACTORS)


def extract(filename: str, content: bytes) -> str:
    """Return the plain text of a document, or raise if it cannot be read."""
    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = _EXTRACTORS.get(extension)
    if extractor is None:
        raise UnsupportedFile(
            f"Unsupported file type '.{extension}'. Supported: {', '.join('.' + e for e in SUPPORTED_EXTENSIONS)}."
        )
    try:
        text = extractor(content).strip()
    except (UnsupportedFile, ExtractionError):
        raise
    except Exception as exc:  # narrow library errors into one type for the caller
        raise ExtractionError(f"Could not read {extension} file: {exc}") from exc
    if not text:
        raise ExtractionError("No readable text was found in the file.")
    return text
